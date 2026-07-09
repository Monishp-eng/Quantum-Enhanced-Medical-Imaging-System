"""
Generate a professional PowerPoint presentation for QT-2.21 
Quantum-Enhanced Medical Imaging for Cancer Detection.
"""
import os
import sys
sys.path.insert(0, '.')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ============================================================
# COLOR PALETTE
# ============================================================
BG_DARK       = RGBColor(0x0A, 0x0C, 0x10)
BG_CARD       = RGBColor(0x12, 0x16, 0x21)
ACCENT_PURPLE = RGBColor(0x7B, 0x61, 0xFF)
ACCENT_BLUE   = RGBColor(0x3B, 0x82, 0xF6)
SUCCESS_GREEN = RGBColor(0x10, 0xB9, 0x81)
WARNING_AMBER = RGBColor(0xF5, 0x9E, 0x0B)
DANGER_RED    = RGBColor(0xF4, 0x3F, 0x5E)
TEXT_PRIMARY  = RGBColor(0xF3, 0xF4, 0xF6)
TEXT_SECONDARY= RGBColor(0x9C, 0xA3, 0xAF)
TEXT_MUTED    = RGBColor(0x6B, 0x72, 0x80)
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
BLACK         = RGBColor(0x00, 0x00, 0x00)

SLIDE_WIDTH  = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT

# ============================================================
# HELPER FUNCTIONS
# ============================================================

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=TEXT_PRIMARY, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_frame(slide, left, top, width, height, items, font_size=16,
                     color=TEXT_SECONDARY, bullet_color=ACCENT_PURPLE):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.space_after = Pt(6)
        p.level = 0
    return txBox

def add_gradient_bar(slide, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    fill = shape.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = ACCENT_PURPLE
    fill.gradient_stops[0].position = 0.0
    fill.gradient_stops[1].color.rgb = ACCENT_BLUE
    fill.gradient_stops[1].position = 1.0
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def slide_header(slide, title, subtitle=None):
    """Add a consistent header bar to each slide."""
    add_gradient_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06))
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
                 title, font_size=32, color=WHITE, bold=True, font_name='Calibri')
    if subtitle:
        add_text_box(slide, Inches(0.8), Inches(0.85), Inches(11), Inches(0.4),
                     subtitle, font_size=16, color=TEXT_SECONDARY, font_name='Calibri')
    # Slide number placeholder area
    add_text_box(slide, Inches(12.0), Inches(7.0), Inches(1.2), Inches(0.4),
                 '', font_size=10, color=TEXT_MUTED, alignment=PP_ALIGN.RIGHT)

# ============================================================
# SLIDE 1 — TITLE SLIDE
# ============================================================
s1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
set_slide_bg(s1, BG_DARK)

# Large gradient accent bar
add_gradient_bar(s1, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.12))

# Title block
add_text_box(s1, Inches(1.5), Inches(1.8), Inches(10), Inches(1.0),
             'Quantum-Enhanced Medical Imaging', font_size=44, color=WHITE, bold=True)
add_text_box(s1, Inches(1.5), Inches(2.8), Inches(10), Inches(0.7),
             'for Cancer Detection', font_size=40, color=ACCENT_PURPLE, bold=True)

# Subtitle
add_text_box(s1, Inches(1.5), Inches(3.8), Inches(10), Inches(0.5),
             'Brain Tumor Classification from MRI Scans using QIEO Optimization',
             font_size=20, color=TEXT_SECONDARY)

# Divider
add_gradient_bar(s1, Inches(1.5), Inches(4.6), Inches(3), Inches(0.04))

# Team & metadata
add_text_box(s1, Inches(1.5), Inches(5.0), Inches(10), Inches(0.4),
             'Project Code: QT-2.21  |  Configuration: 6 Models  |  Dataset: 7,023 MRI Scans',
             font_size=14, color=TEXT_MUTED)

add_text_box(s1, Inches(1.5), Inches(5.5), Inches(10), Inches(0.4),
             'Team: Loki (Data Engineering) · SK (Feature Engineering) · Monish (Quantum Optimization) · Akshya (Evaluation & Explainability)',
             font_size=14, color=TEXT_MUTED)

# ============================================================
# SLIDE 2 — PROBLEM STATEMENT & OBJECTIVES
# ============================================================
s2 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s2, BG_DARK)
slide_header(s2, 'Problem Statement & Objectives')

# Problem statement card
add_rect(s2, Inches(0.6), Inches(1.5), Inches(5.8), Inches(2.6), BG_CARD, RGBColor(0x1E, 0x22, 0x33))
add_text_box(s2, Inches(0.9), Inches(1.6), Inches(5.4), Inches(0.4),
             '🧠 Clinical Problem', font_size=18, color=ACCENT_PURPLE, bold=True)
add_bullet_frame(s2, Inches(0.9), Inches(2.1), Inches(5.4), Inches(1.8), [
    '• Brain tumors require fast, accurate, non-invasive diagnosis',
    '• Manual MRI analysis is time-consuming & error-prone',
    '• Inter-observer variability between radiologists is high',
    '• Early detection directly impacts patient survival rates',
], font_size=15, color=TEXT_SECONDARY)

# Objectives card
add_rect(s2, Inches(6.8), Inches(1.5), Inches(5.8), Inches(2.6), BG_CARD, RGBColor(0x1E, 0x22, 0x33))
add_text_box(s2, Inches(7.1), Inches(1.6), Inches(5.4), Inches(0.4),
             '🎯 Project Objectives', font_size=18, color=SUCCESS_GREEN, bold=True)
add_bullet_frame(s2, Inches(7.1), Inches(2.1), Inches(5.4), Inches(1.8), [
    '① Classify MRI scans: Glioma / Meningioma / Pituitary / Healthy',
    '② Extract features via preprocessing & feature engineering',
    '③ Optimize using Quantum-Inspired Evolutionary Algorithm',
    '④ Maximize sensitivity for early cancer detection',
    '⑤ Compare quantum vs. classical ML pipelines statistically',
], font_size=14, color=TEXT_SECONDARY)

# Dataset card
add_rect(s2, Inches(0.6), Inches(4.5), Inches(12.0), Inches(2.5), BG_CARD, RGBColor(0x1E, 0x22, 0x33))
add_text_box(s2, Inches(0.9), Inches(4.6), Inches(11.4), Inches(0.4),
             '📊 Dataset Overview', font_size=18, color=WARNING_AMBER, bold=True)

# 4 class mini cards
classes = [
    ('Glioma', '1,321 images', 'Malignant tumor of glial cells', DANGER_RED),
    ('Meningioma', '1,339 images', 'Mostly benign, from meninges', WARNING_AMBER),
    ('Pituitary', '1,457 images', 'Benign pituitary adenoma', ACCENT_PURPLE),
    ('No Tumor', '1,595 images', 'Healthy brain tissue', SUCCESS_GREEN),
]
for i, (name, count, desc, accent) in enumerate(classes):
    x = Inches(0.9 + i * 2.9)
    add_rect(s2, x, Inches(5.2), Inches(2.6), Inches(1.5), RGBColor(0x0A, 0x0E, 0x18), accent)
    add_text_box(s2, x + Inches(0.15), Inches(5.3), Inches(2.3), Inches(0.35),
                 name, font_size=16, color=accent, bold=True)
    add_text_box(s2, x + Inches(0.15), Inches(5.65), Inches(2.3), Inches(0.3),
                 count, font_size=13, color=WHITE, bold=True)
    add_text_box(s2, x + Inches(0.15), Inches(5.95), Inches(2.3), Inches(0.5),
                 desc, font_size=11, color=TEXT_MUTED)

# ============================================================
# SLIDE 3 — SYSTEM ARCHITECTURE
# ============================================================
s3 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s3, BG_DARK)
slide_header(s3, 'System Architecture & Pipeline')

# Pipeline flow boxes
pipeline_steps = [
    ('1. Raw MRI Input', 'Brain MRI scans\n(224×224 pixels)', ACCENT_BLUE),
    ('2. Preprocessing', 'Skull Strip → CLAHE\n→ Denoise → Normalize', ACCENT_PURPLE),
    ('3. Feature Extraction', 'ResNet-50 (2048-D)\n+ Handcrafted (98-D)\n= 2146-D Fused Vector', WARNING_AMBER),
    ('4. Optimization', 'QIEO Feature Selection\n+ QIEO Hyperparam\nSearch (56-bit)', SUCCESS_GREEN),
    ('5. Classification', 'SVM / XGBoost /\nResNet-50 CNN\n→ 4-class Output', DANGER_RED),
    ('6. Explainability', 'Grad-CAM Heatmap\n+ Saliency Maps\n+ Clinical Reasoning', ACCENT_BLUE),
]
for i, (title, desc, accent) in enumerate(pipeline_steps):
    x = Inches(0.4 + i * 2.1)
    add_rect(s3, x, Inches(1.8), Inches(1.9), Inches(2.4), BG_CARD, accent)
    add_text_box(s3, x + Inches(0.1), Inches(1.9), Inches(1.7), Inches(0.4),
                 title, font_size=13, color=accent, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(s3, x + Inches(0.1), Inches(2.4), Inches(1.7), Inches(1.5),
                 desc, font_size=11, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)
    # Arrow between boxes
    if i < len(pipeline_steps) - 1:
        add_text_box(s3, x + Inches(1.9), Inches(2.7), Inches(0.3), Inches(0.4),
                     '→', font_size=20, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

# Bottom detail cards
add_rect(s3, Inches(0.6), Inches(4.8), Inches(5.8), Inches(2.3), BG_CARD, RGBColor(0x1E, 0x22, 0x33))
add_text_box(s3, Inches(0.9), Inches(4.9), Inches(5.4), Inches(0.35),
             '🔬 Preprocessing Pipeline Details', font_size=16, color=ACCENT_PURPLE, bold=True)
add_bullet_frame(s3, Inches(0.9), Inches(5.3), Inches(5.4), Inches(1.6), [
    '• Skull Stripping: Otsu thresholding + contour masking',
    '• CLAHE: clip_limit=2.0, tile_grid=(8×8)',
    '• Gaussian Blur: kernel=(3×3) for noise reduction',
    '• Min-Max Normalization: scale to [0.0, 1.0]',
    '• Stratified Split: 70% train / 15% val / 15% test',
], font_size=13, color=TEXT_SECONDARY)

add_rect(s3, Inches(6.8), Inches(4.8), Inches(5.8), Inches(2.3), BG_CARD, RGBColor(0x1E, 0x22, 0x33))
add_text_box(s3, Inches(7.1), Inches(4.9), Inches(5.4), Inches(0.35),
             '🧬 Feature Engineering (2,146-D)', font_size=16, color=WARNING_AMBER, bold=True)
add_bullet_frame(s3, Inches(7.1), Inches(5.3), Inches(5.4), Inches(1.6), [
    '• Deep CNN: ResNet-50 penultimate layer → 2,048 dims',
    '• GLCM: Contrast, Correlation, Homogeneity → 20 dims',
    '• HOG: Gradient orientation histograms → 36 dims',
    '• DWT: 4-level Haar wavelet subbands → 32 dims',
    '• Intensity: Mean, Std, Skewness, Kurtosis → 10 dims',
], font_size=13, color=TEXT_SECONDARY)

# ============================================================
# SLIDE 4 — QIEO ALGORITHM
# ============================================================
s4 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s4, BG_DARK)
slide_header(s4, 'Quantum-Inspired Evolutionary Optimization (QIEO)',
             'Core Innovation: Q-bit Superposition + Quantum Rotation Gates')

add_rect(s4, Inches(0.6), Inches(1.5), Inches(5.8), Inches(5.3), BG_CARD, RGBColor(0x1E, 0x22, 0x33))
add_text_box(s4, Inches(0.9), Inches(1.6), Inches(5.4), Inches(0.4),
             '⚛️ Mathematical Formulation', font_size=18, color=ACCENT_PURPLE, bold=True)
add_bullet_frame(s4, Inches(0.9), Inches(2.1), Inches(5.4), Inches(4.5), [
    'Each feature/parameter is a Q-bit state vector:',
    '    q_j = [α_j, β_j]ᵀ  where |α|² + |β|² = 1',
    '',
    '• |α_j|² = probability of collapse to 0 (unselected)',
    '• |β_j|² = probability of collapse to 1 (selected)',
    '',
    'Quantum Rotation Gate Update:',
    '    α_j\' = α_j·cos(θ) − β_j·sin(θ)',
    '    β_j\' = α_j·sin(θ) + β_j·cos(θ)',
    '',
    'θ is dynamically computed using a lookup table',
    'to guide Q-bits towards globally best chromosome.',
    '',
    'Initial state: Equal superposition (α = β = 1/√2)',
], font_size=13, color=TEXT_SECONDARY)

add_rect(s4, Inches(6.8), Inches(1.5), Inches(5.8), Inches(2.4), BG_CARD, RGBColor(0x1E, 0x22, 0x33))
add_text_box(s4, Inches(7.1), Inches(1.6), Inches(5.4), Inches(0.4),
             '🎯 QIEO Feature Selection', font_size=18, color=SUCCESS_GREEN, bold=True)
add_bullet_frame(s4, Inches(7.1), Inches(2.1), Inches(5.4), Inches(1.6), [
    '• Search Space: 2,146 binary Q-bits (one per feature)',
    '• Population: 10 chromosomes × 15 generations',
    '• Fitness: CV_Accuracy − 0.01 × (selected / total)',
    '• Result: Selected 1,072 optimal features (50% compression)',
    '• 3-fold stratified CV on 1,000-sample subset for speed',
], font_size=13, color=TEXT_SECONDARY)

add_rect(s4, Inches(6.8), Inches(4.2), Inches(5.8), Inches(2.6), BG_CARD, RGBColor(0x1E, 0x22, 0x33))
add_text_box(s4, Inches(7.1), Inches(4.3), Inches(5.4), Inches(0.4),
             '⚙️ QIEO Hyperparameter Optimization', font_size=18, color=WARNING_AMBER, bold=True)
add_bullet_frame(s4, Inches(7.1), Inches(4.8), Inches(5.4), Inches(1.8), [
    '• 7 XGBoost params × 8 bits = 56-bit chromosome',
    '• n_estimators: 50–550, max_depth: 3–15',
    '• learning_rate: 0.01–0.30, subsample: 0.5–1.0',
    '• Population: 8 individuals × 10 generations',
    '• Binary-to-float decoder maps bits → continuous space',
    '• 3-fold CV on 1,500-sample stratified subset',
], font_size=13, color=TEXT_SECONDARY)

# ============================================================
# SLIDE 5 — MODEL CONFIGURATIONS
# ============================================================
s5 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s5, BG_DARK)
slide_header(s5, '6 Model Configurations: Classical vs. Quantum')

configs = [
    ('Baseline A', 'Mutual Info\n(Top 100)', 'GridSearchCV', 'SVM (RBF)', '~20 min', TEXT_SECONDARY),
    ('Baseline B', 'SelectKBest\n(Top 100)', 'Optuna TPE\n(25 trials)', 'XGBoost', '~30 min', TEXT_SECONDARY),
    ('Baseline C', 'End-to-end\nCNN Images', 'Transfer\nLearning', 'ResNet-50\n(5 epochs)', '~3–4 hrs', TEXT_SECONDARY),
    ('Quantum A', 'QIEO Mask\n(1072 feats)', 'GridSearchCV', 'SVM (RBF)', '~45 min', ACCENT_PURPLE),
    ('Quantum B', 'QIEO Mask\n(1072 feats)', 'QIEO 56-bit\nBinary', 'XGBoost', '~35 min', ACCENT_PURPLE),
    ('Quantum C', 'QIEO\nInformed', 'QIEO\nInformed', 'ResNet-50', 'Shared', ACCENT_PURPLE),
]

headers = ['Config', 'Feature Selection', 'Hyperparam Search', 'Classifier', 'Training Time']
col_widths = [Inches(2.0), Inches(2.2), Inches(2.2), Inches(2.2), Inches(1.8)]
table_left = Inches(0.8)
table_top = Inches(1.8)

for j, h in enumerate(headers):
    x = table_left + sum(col_widths[:j], Emu(0))
    add_rect(s5, x, table_top, col_widths[j] - Emu(Pt(2)), Inches(0.5), ACCENT_PURPLE)
    add_text_box(s5, x + Inches(0.1), table_top + Inches(0.05), col_widths[j] - Inches(0.2), Inches(0.4),
                 h, font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

for i, (name, feat, hyper, clf, time, accent) in enumerate(configs):
    y = table_top + Inches(0.55 + i * 0.85)
    row_bg = RGBColor(0x14, 0x18, 0x26) if accent == ACCENT_PURPLE else BG_CARD
    vals = [name, feat, hyper, clf, time]
    for j, val in enumerate(vals):
        x = table_left + sum(col_widths[:j], Emu(0))
        add_rect(s5, x, y, col_widths[j] - Emu(Pt(2)), Inches(0.8), row_bg, RGBColor(0x1E, 0x22, 0x33))
        c = accent if j == 0 else TEXT_SECONDARY
        add_text_box(s5, x + Inches(0.1), y + Inches(0.1), col_widths[j] - Inches(0.2), Inches(0.6),
                     val, font_size=12, color=c, bold=(j==0), alignment=PP_ALIGN.CENTER)

# Note at bottom
add_text_box(s5, Inches(0.8), Inches(7.0), Inches(11), Inches(0.35),
             '⚛️ Purple rows = Quantum-optimized configurations using QIEO solver',
             font_size=12, color=ACCENT_PURPLE)

# ============================================================
# SLIDE 6 — RESULTS TABLE
# ============================================================
s6 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s6, BG_DARK)
slide_header(s6, 'Performance Evaluation Results', 'Test Set: N = 1,054 samples across 4 classes')

results_headers = ['Config', 'Accuracy', 'Sensitivity', 'Specificity', 'Precision', 'F1-Macro', 'AUC-ROC']
results_data = [
    ('Baseline A', '87.00%', '87.00%', '95.67%', '87.48%', '86.99%', '97.59%'),
    ('Baseline B', '87.00%', '87.00%', '95.67%', '87.17%', '86.96%', '97.46%'),
    ('Baseline C', '95.33%', '95.33%', '98.44%', '95.46%', '95.34%', '99.41%'),
    ('Quantum A', '89.84%', '89.84%', '96.61%', '89.98%', '89.83%', '98.62%'),
    ('Quantum B', '89.29%', '89.29%', '96.43%', '89.50%', '89.28%', '97.99%'),
    ('Quantum C', '95.33%', '95.33%', '98.44%', '95.46%', '95.34%', '99.41%'),
]

r_col_widths = [Inches(1.8), Inches(1.5), Inches(1.6), Inches(1.5), Inches(1.5), Inches(1.5), Inches(1.5)]
r_left = Inches(0.5)
r_top = Inches(1.8)

for j, h in enumerate(results_headers):
    x = r_left + sum(r_col_widths[:j], Emu(0))
    add_rect(s6, x, r_top, r_col_widths[j] - Emu(Pt(2)), Inches(0.5), ACCENT_PURPLE)
    add_text_box(s6, x + Inches(0.05), r_top + Inches(0.05), r_col_widths[j] - Inches(0.1), Inches(0.4),
                 h, font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

for i, row in enumerate(results_data):
    y = r_top + Inches(0.55 + i * 0.65)
    is_quantum = 'Quantum' in row[0]
    row_bg = RGBColor(0x14, 0x18, 0x26) if is_quantum else BG_CARD
    for j, val in enumerate(row):
        x = r_left + sum(r_col_widths[:j], Emu(0))
        add_rect(s6, x, y, r_col_widths[j] - Emu(Pt(2)), Inches(0.6), row_bg, RGBColor(0x1E, 0x22, 0x33))
        c = ACCENT_PURPLE if (is_quantum and j == 0) else (SUCCESS_GREEN if (row[0] in ('Baseline C', 'Quantum C') and j > 0) else TEXT_PRIMARY)
        add_text_box(s6, x + Inches(0.05), y + Inches(0.1), r_col_widths[j] - Inches(0.1), Inches(0.4),
                     val, font_size=13, color=c, bold=(j==0 or (row[0] in ('Baseline C', 'Quantum C') and j > 0)),
                     alignment=PP_ALIGN.CENTER)

# Improvement badges
add_rect(s6, Inches(0.6), Inches(5.8), Inches(5.8), Inches(1.4), BG_CARD, SUCCESS_GREEN)
add_text_box(s6, Inches(0.9), Inches(5.9), Inches(5.4), Inches(0.35),
             '📈 Quantum Improvement over Classical', font_size=16, color=SUCCESS_GREEN, bold=True)
add_bullet_frame(s6, Inches(0.9), Inches(6.3), Inches(5.4), Inches(0.8), [
    '• Quantum A vs Baseline A:  +3.26% Accuracy  |  +3.27% F1-Score',
    '• Quantum B vs Baseline B:  +2.63% Accuracy  |  +2.66% F1-Score',
], font_size=14, color=TEXT_PRIMARY)

# Statistical significance
add_rect(s6, Inches(6.8), Inches(5.8), Inches(5.8), Inches(1.4), BG_CARD, ACCENT_PURPLE)
add_text_box(s6, Inches(7.1), Inches(5.9), Inches(5.4), Inches(0.35),
             '⚖️ Statistical Significance (McNemar\'s Test)', font_size=16, color=ACCENT_PURPLE, bold=True)
add_bullet_frame(s6, Inches(7.1), Inches(6.3), Inches(5.4), Inches(0.8), [
    '• Baseline A vs Quantum A:  p = 0.0037  ✅ Significant',
    '• Baseline B vs Quantum B:  p = 0.0110  ✅ Significant',
], font_size=14, color=TEXT_PRIMARY)

# ============================================================
# SLIDE 7 — PER-CLASS SENSITIVITY & AUC-ROC
# ============================================================
s7 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s7, BG_DARK)
slide_header(s7, 'Per-Class Sensitivity & AUC-ROC (Quantum A)', 'Maximizing Sensitivity for Early Cancer Detection')

# Sensitivity cards
sens_data = [
    ('Glioma', '82.78%', DANGER_RED),
    ('Meningioma', '86.81%', WARNING_AMBER),
    ('Pituitary', '95.24%', ACCENT_PURPLE),
    ('No Tumor', '94.51%', SUCCESS_GREEN),
]
add_text_box(s7, Inches(0.8), Inches(1.5), Inches(5), Inches(0.4),
             '🎯 Per-Class Sensitivity (True Positive Rate)', font_size=18, color=WHITE, bold=True)
for i, (name, val, accent) in enumerate(sens_data):
    x = Inches(0.8 + i * 3.0)
    add_rect(s7, x, Inches(2.0), Inches(2.7), Inches(1.8), BG_CARD, accent)
    add_text_box(s7, x + Inches(0.2), Inches(2.1), Inches(2.3), Inches(0.35),
                 name, font_size=14, color=accent, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(s7, x + Inches(0.2), Inches(2.6), Inches(2.3), Inches(0.8),
                 val, font_size=32, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# AUC-ROC cards
auc_data = [
    ('Glioma', '97.79%', DANGER_RED),
    ('Meningioma', '97.54%', WARNING_AMBER),
    ('Pituitary', '99.52%', ACCENT_PURPLE),
    ('No Tumor', '99.62%', SUCCESS_GREEN),
]
add_text_box(s7, Inches(0.8), Inches(4.2), Inches(5), Inches(0.4),
             '📊 Per-Class AUC-ROC (One-vs-Rest)', font_size=18, color=WHITE, bold=True)
for i, (name, val, accent) in enumerate(auc_data):
    x = Inches(0.8 + i * 3.0)
    add_rect(s7, x, Inches(4.7), Inches(2.7), Inches(1.8), BG_CARD, accent)
    add_text_box(s7, x + Inches(0.2), Inches(4.8), Inches(2.3), Inches(0.35),
                 name, font_size=14, color=accent, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(s7, x + Inches(0.2), Inches(5.3), Inches(2.3), Inches(0.8),
                 val, font_size=32, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 8 — EXPLAINABILITY
# ============================================================
s8 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s8, BG_DARK)
slide_header(s8, 'Model Explainability & Clinical Trust',
             'Visual + Textual Reasoning for Every Diagnosis')

add_rect(s8, Inches(0.6), Inches(1.5), Inches(5.8), Inches(5.3), BG_CARD, RGBColor(0x1E, 0x22, 0x33))
add_text_box(s8, Inches(0.9), Inches(1.6), Inches(5.4), Inches(0.4),
             '🔥 Grad-CAM Heatmap', font_size=18, color=DANGER_RED, bold=True)
add_bullet_frame(s8, Inches(0.9), Inches(2.1), Inches(5.4), Inches(2.0), [
    '• Computed from ResNet-50 layer4 (base.7)',
    '• Gradient-weighted class activation mapping',
    '• L_cam = ReLU(Σ α_k · A_k)  (weighted activations)',
    '• Overlaid on MRI scan for spatial visualization',
    '• Shows WHERE the CNN focused its attention',
    '',
    '🔬 Vanilla Saliency Maps',
    '• Absolute gradient of class score w.r.t. input pixels',
    '• S(x) = max_c |∂S_c / ∂x|',
    '• Highlights micro-level edges and boundaries',
], font_size=14, color=TEXT_SECONDARY)

add_rect(s8, Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.3), BG_CARD, RGBColor(0x1E, 0x22, 0x33))
add_text_box(s8, Inches(7.1), Inches(1.6), Inches(5.4), Inches(0.4),
             '🩺 Clinical Reasoning Engine (NEW)', font_size=18, color=SUCCESS_GREEN, bold=True)
add_bullet_frame(s8, Inches(7.1), Inches(2.1), Inches(5.4), Inches(4.5), [
    'For every prediction, the system generates a detailed',
    'clinical-style explanation including:',
    '',
    '• GLCM Texture Analysis: contrast, homogeneity, energy',
    '• Intensity Profile: mean brightness, variance, skewness',
    '• Wavelet Energy: multi-scale frequency response',
    '• HOG Gradient Structure: edge consistency analysis',
    '',
    '• Grad-CAM Spatial Focus: identifies exact brain region',
    '  (e.g., "superior left region, highly localized")',
    '',
    '• Differential Diagnosis: runner-up class + confidence',
    '  margin analysis with radiologist review warnings',
], font_size=13, color=TEXT_SECONDARY)

# ============================================================
# SLIDE 9 — LIVE DEMO DASHBOARD
# ============================================================
s9 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s9, BG_DARK)
slide_header(s9, 'Live Demo Dashboard', 'Interactive Web Application at http://localhost:8000')

features = [
    ('🧠 MRI Upload & Preview', 'Drag-and-drop or browse\nfor any MRI scan.\nLoad random test scan.'),
    ('⚛️ Model Selection', 'Switch between 4 models:\nQuantum A/B, Baseline A/B.\nInstant re-prediction.'),
    ('📊 Probability Distribution', 'Real-time class confidence\nbars for all 4 tumor types\nwith predicted diagnosis.'),
    ('🔥 Grad-CAM + Saliency', 'Side-by-side explainability\nmaps generated on-the-fly\nfrom ResNet-50 layer4.'),
    ('🩺 Diagnostic Explanation', 'Feature-level clinical\nreasoning: GLCM, HOG,\nwavelet, Grad-CAM region.'),
    ('📈 Evaluation Report', 'Dedicated /evaluation page\nwith confusion matrices,\nall metrics, McNemar stats.'),
]
for i, (title, desc) in enumerate(features):
    row = i // 3
    col = i % 3
    x = Inches(0.6 + col * 4.2)
    y = Inches(1.8 + row * 2.8)
    add_rect(s9, x, y, Inches(3.8), Inches(2.4), BG_CARD, RGBColor(0x1E, 0x22, 0x33))
    add_text_box(s9, x + Inches(0.2), y + Inches(0.15), Inches(3.4), Inches(0.4),
                 title, font_size=16, color=ACCENT_PURPLE, bold=True)
    add_text_box(s9, x + Inches(0.2), y + Inches(0.65), Inches(3.4), Inches(1.5),
                 desc, font_size=13, color=TEXT_SECONDARY)

# ============================================================
# SLIDE 10 — CLINICAL LIMITATIONS & FUTURE WORK
# ============================================================
s10 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s10, BG_DARK)
slide_header(s10, 'Clinical Limitations & Future Directions')

add_rect(s10, Inches(0.6), Inches(1.5), Inches(5.8), Inches(5.3), BG_CARD, RGBColor(0x1E, 0x22, 0x33))
add_text_box(s10, Inches(0.9), Inches(1.6), Inches(5.4), Inches(0.4),
             '⚠️ Current Limitations', font_size=18, color=WARNING_AMBER, bold=True)
add_bullet_frame(s10, Inches(0.9), Inches(2.2), Inches(5.4), Inches(4.3), [
    '1. Resolution Downsampling',
    '   Images resized to 224×224 pixels. May lose',
    '   micro-metastases or tiny lesions.',
    '',
    '2. Single Modality (MRI Only)',
    '   Trained exclusively on MRI scans. Cannot process',
    '   CT scans due to different intensity profiles.',
    '',
    '3. Single Dataset Bias',
    '   Trained on Kaggle Brain Tumor MRI dataset only.',
    '   Generalization to other scanners (1.5T vs 3T)',
    '   and institutions is unproven.',
    '',
    '4. Grad-CAM Resolution',
    '   Heatmaps from 7×7 spatial grid do not perfectly',
    '   align with microscopic tumor boundaries.',
], font_size=13, color=TEXT_SECONDARY)

add_rect(s10, Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.3), BG_CARD, RGBColor(0x1E, 0x22, 0x33))
add_text_box(s10, Inches(7.1), Inches(1.6), Inches(5.4), Inches(0.4),
             '🚀 Future Directions', font_size=18, color=SUCCESS_GREEN, bold=True)
add_bullet_frame(s10, Inches(7.1), Inches(2.2), Inches(5.4), Inches(4.3), [
    '1. 3D Volumetric Segmentation',
    '   Expand preprocessing and CNN modules to ingest',
    '   entire 3D MRI voxel volumes (3D-UNet, V-Net).',
    '',
    '2. True Quantum Hardware',
    '   Migrate QIEO to NISQ quantum computers using',
    '   Qiskit or Pennylane quantum circuits.',
    '',
    '3. Multi-Modal Clinical Fusion',
    '   Fuse MRI outputs with clinical text (radiology',
    '   reports) and genomic biomarker data.',
    '',
    '4. Multi-Institutional Validation',
    '   Validate on diverse hospital datasets to ensure',
    '   cross-scanner and cross-population robustness.',
    '',
    '5. Federated Learning',
    '   Enable privacy-preserving model training across',
    '   hospitals without sharing raw patient data.',
], font_size=13, color=TEXT_SECONDARY)

# ============================================================
# SLIDE 11 — CONCLUSION & THANK YOU
# ============================================================
s11 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s11, BG_DARK)
add_gradient_bar(s11, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.12))

add_text_box(s11, Inches(1.5), Inches(1.5), Inches(10), Inches(0.8),
             'Key Takeaways', font_size=36, color=WHITE, bold=True)

takeaways = [
    ('✅', '89.84% Accuracy with Quantum A (QIEO + SVM) — +3.26% over classical baseline'),
    ('✅', '95.33% Accuracy with ResNet-50 CNN — highest overall performance'),
    ('✅', 'p = 0.0037 — Quantum optimization is statistically significant'),
    ('✅', 'Full explainability: Grad-CAM + Saliency + Clinical Reasoning Engine'),
    ('✅', 'End-to-end system: Preprocessing → Feature Fusion → QIEO → Classification → Dashboard'),
]
for i, (icon, text) in enumerate(takeaways):
    y = Inches(2.5 + i * 0.7)
    add_rect(s11, Inches(1.5), y, Inches(10), Inches(0.6), BG_CARD, RGBColor(0x1E, 0x22, 0x33))
    add_text_box(s11, Inches(1.7), y + Inches(0.05), Inches(9.6), Inches(0.5),
                 f'{icon}  {text}', font_size=16, color=TEXT_PRIMARY, bold=False)

add_gradient_bar(s11, Inches(4), Inches(6.0), Inches(5), Inches(0.04))

add_text_box(s11, Inches(1.5), Inches(6.3), Inches(10), Inches(0.6),
             'Thank You!', font_size=36, color=ACCENT_PURPLE, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(s11, Inches(1.5), Inches(6.9), Inches(10), Inches(0.4),
             'QT-2.21 Diagnostics  •  Questions?',
             font_size=16, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)

# ============================================================
# SAVE
# ============================================================
output_path = os.path.join('reports', 'QT221_Quantum_Enhanced_Medical_Imaging.pptx')
os.makedirs('reports', exist_ok=True)
prs.save(output_path)
print(f'\nPresentation saved to: {output_path}')
print(f'Total slides: {len(prs.slides)}')
