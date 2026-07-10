"""
Generate a professional, high-scoring hackathon presentation for the RIT Quant-A-Than 2026 
using the exact slide template and graphics from 'RIT - Quantathan 2026 - PPT.pptx'.
"""
import os
import sys
import copy
import io
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ============================================================
# COLOR PALETTE (For White-Theme Content Slides)
# ============================================================
COLOR_DARK      = RGBColor(0x06, 0x08, 0x14)  # Main text
COLOR_MUTED     = RGBColor(0x4B, 0x55, 0x63)  # Secondary text
COLOR_ACCENT    = RGBColor(0x7B, 0x61, 0xFF)  # RIT Purple Accent
COLOR_BLUE      = RGBColor(0x3B, 0x82, 0xF6)  # Accent Blue
COLOR_GREEN     = RGBColor(0x10, 0xB9, 0x81)  # Success Green
COLOR_RED       = RGBColor(0xEF, 0x44, 0x44)  # Danger Red
COLOR_AMBER     = RGBColor(0xF5, 0x9E, 0x0B)  # Warning Amber
COLOR_CARD_BG   = RGBColor(0xF9, 0xFA, 0xFB)  # Light gray card fill
COLOR_CARD_BORD = RGBColor(0xE5, 0xE7, 0xEB)  # Card border

# File paths
TEMPLATE_PATH = r"C:\Users\Monish P\Downloads\RIT - Quantathan 2026 - PPT.pptx"
OUTPUT_PATH   = r"C:\Users\Monish P\Downloads\RIT_Quantathan_2026_Final.pptx"

def main():
    if not os.path.exists(TEMPLATE_PATH):
        print(f"Error: Template file not found at {TEMPLATE_PATH}")
        sys.exit(1)
        
    prs = Presentation(TEMPLATE_PATH)
    
    # Template base slides:
    # prs.slides[0] is the Title Slide (Dark Theme)
    # prs.slides[1] is the Content Slide Template (White Theme)
    s1 = prs.slides[0]
    s2 = prs.slides[1]
    
    # ============================================================
    # SLIDE 1 — TITLE SLIDE (Modify in place)
    # ============================================================
    # Replace placeholder text for candidate and organization
    for shape in s1.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text
            if "Name of the Candidate" in text:
                shape.text_frame.paragraphs[0].runs[0].text = "Presented by: Monish P"
            elif "Institution / Organization" in text:
                shape.text_frame.paragraphs[0].runs[0].text = "Rajalakshmi Institute of Technology"
                
    # Add project title and subtitle in the middle-left empty area of slide 1
    # Left=0.6, Top=4.2, Width=7.5, Height=1.5
    title_box = s1.shapes.add_textbox(Inches(0.6), Inches(4.0), Inches(7.5), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    p1.text = "QT-2.21: Quantum-Enhanced Brain Tumor Classification"
    p1.font.size = Pt(28)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p1.font.name = 'Calibri'
    p1.space_after = Pt(6)
    
    p2 = tf.add_paragraph()
    p2.text = "Real-time Clinical Diagnosis Pipeline using Deep Feature Fusion & QIEO Solver"
    p2.font.size = Pt(14)
    p2.font.color.rgb = COLOR_BLUE
    p2.font.name = 'Calibri'
    
    print("Modified Slide 1 (Title slide)")

    # ============================================================
    # HELPER TO CLONE CONTENT SLIDE (Slide 2 layout)
    # ============================================================
    def create_cloned_slide():
        blank_layout = prs.slide_layouts[6]
        new_slide = prs.slides.add_slide(blank_layout)
        
        # 1. Copy slide background image & relationship
        try:
            rId = s2.background.fill._xPr.xpath('.//a:blip/@r:embed')[0]
            rel = s2.part.rels[rId]
            img_part = s2.part.related_part(rId)
            new_rId = new_slide.part.relate_to(img_part, rel.reltype)
            
            new_bg = copy.deepcopy(s2.background._element)
            new_bg.xpath('.//a:blip')[0].set('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed', new_rId)
            new_slide.background._element.clear()
            for child in new_bg:
                new_slide.background._element.append(child)
        except Exception as e:
            print("Background copy warning:", e)
            
        # 2. Copy shapes and their relationships
        for shape in s2.shapes:
            new_shape_el = copy.deepcopy(shape.element)
            blips = new_shape_el.xpath('.//a:blip')
            if blips:
                for blip in blips:
                    rId = blip.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
                    if rId:
                        try:
                            img_part = s2.part.related_part(rId)
                            rel = s2.part.rels[rId]
                            new_rId = new_slide.part.relate_to(img_part, rel.reltype)
                            blip.set('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed', new_rId)
                        except Exception as e:
                            print("Shape image relationship copy error:", e)
            new_slide.shapes._spTree.append(new_shape_el)
            
        return new_slide

    # Helper functions for slide content
    def add_slide_title(slide, title_text):
        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(9.2), Inches(0.5))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = COLOR_DARK
        p.font.name = 'Calibri'
        
    def add_card(slide, left, top, width, height, bg_color=COLOR_CARD_BG, border_color=COLOR_CARD_BORD):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        if border_color:
            shape.line.color.rgb = border_color
            shape.line.width = Pt(1.5)
        else:
            shape.line.fill.background()
        shape.shadow.inherit = False
        return shape
        
    def add_textbox(slide, left, top, width, height, text, font_size=14, color=COLOR_DARK, bold=False, alignment=PP_ALIGN.LEFT):
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.alignment = alignment
        return box

    def add_bullet_list(slide, left, top, width, height, items, font_size=13, color=COLOR_DARK):
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(font_size)
            p.font.color.rgb = color
            p.font.name = 'Calibri'
            p.space_after = Pt(4)
            p.level = 0
        return box

    # ============================================================
    # SLIDE 2 — PROBLEM STATEMENT & OBJECTIVES
    # ============================================================
    s_prob = create_cloned_slide()
    add_slide_title(s_prob, "Problem Statement & Objectives")
    
    # Left Column Card (Clinical Challenges)
    add_card(s_prob, Inches(0.6), Inches(2.0), Inches(4.5), Inches(4.6))
    add_textbox(s_prob, Inches(0.8), Inches(2.1), Inches(4.1), Inches(0.4), "🧠 Clinical Challenges", font_size=16, color=COLOR_ACCENT, bold=True)
    add_bullet_list(s_prob, Inches(0.8), Inches(2.6), Inches(4.1), Inches(3.8), [
        "• Brain tumor diagnosis from MRI is a critical clinical step.",
        "• Manual radiological review is time-consuming and subjective.",
        "• High inter-observer variability can delay treatment plans.",
        "• Early, highly sensitive detection directly impacts patient survival rates.",
        "• Computer-aided systems are needed to provide diagnostic assurance."
    ], font_size=13, color=COLOR_MUTED)

    # Right Column Card (Project Objectives)
    add_card(s_prob, Inches(5.4), Inches(2.0), Inches(4.5), Inches(4.6))
    add_textbox(s_prob, Inches(5.6), Inches(2.1), Inches(4.1), Inches(0.4), "🎯 Project Objectives", font_size=16, color=COLOR_BLUE, bold=True)
    add_bullet_list(s_prob, Inches(5.6), Inches(2.6), Inches(4.1), Inches(3.8), [
        "① Classify MRI scans: Glioma / Meningioma / Pituitary / Healthy.",
        "② Develop a high-fidelity preprocessing & normalization pipeline.",
        "③ Build a hybrid deep (ResNet50) + handcrafted features fusion stream.",
        "④ Implement Quantum-Inspired Evolutionary Optimization (QIEO).",
        "⑤ Compare quantum-inspired optimized pipelines with classical baselines."
    ], font_size=13, color=COLOR_MUTED)

    print("Created Slide 2: Problem & Objectives")

    # ============================================================
    # SLIDE 3 — DATASET CHARACTERIZATION
    # ============================================================
    s_data = create_cloned_slide()
    add_slide_title(s_data, "Dataset Characterization & Split")
    
    # Dataset Description Text Box
    add_textbox(s_data, Inches(0.6), Inches(2.0), Inches(9.2), Inches(0.6),
                "Based on the Kaggle Brain Tumor MRI Dataset containing 7,023 images across 4 distinct classes. Splits are stratified (70% Train, 15% Val, 15% Test) to ensure distribution consistency.",
                font_size=13, color=COLOR_MUTED)
                
    # 4 Grid Cards
    classes_data = [
        ("Glioma", "1,321 Scans", "Malignant, aggressive glial tumors", COLOR_RED),
        ("Meningioma", "1,339 Scans", "Extra-axial meninges tumors (benign)", COLOR_AMBER),
        ("Pituitary", "1,457 Scans", "Sellar pituitary adenomas", COLOR_ACCENT),
        ("No Tumor", "1,595 Scans", "Healthy brain parenchyma baseline", COLOR_GREEN),
    ]
    for i, (name, count, desc, color) in enumerate(classes_data):
        row = i // 2
        col = i % 2
        x = Inches(0.6 + col * 4.8)
        y = Inches(2.8 + row * 2.0)
        
        # Subcard
        add_card(s_data, x, y, Inches(4.4), Inches(1.8), bg_color=COLOR_CARD_BG, border_color=color)
        add_textbox(s_data, x + Inches(0.2), y + Inches(0.15), Inches(4.0), Inches(0.35), name, font_size=15, color=color, bold=True)
        add_textbox(s_data, x + Inches(0.2), y + Inches(0.55), Inches(4.0), Inches(0.3), count, font_size=13, color=COLOR_DARK, bold=True)
        add_textbox(s_data, x + Inches(0.2), y + Inches(0.9), Inches(4.0), Inches(0.7), desc, font_size=11, color=COLOR_MUTED)

    print("Created Slide 3: Dataset Characterization")

    # ============================================================
    # SLIDE 4 — SYSTEM ARCHITECTURE
    # ============================================================
    s_arch = create_cloned_slide()
    add_slide_title(s_arch, "System Architecture & Pipeline")
    
    # 5 Horizontal Process Steps
    steps = [
        ("1. MRI Input", "Raw brain scan\n(224×224×3)", COLOR_BLUE),
        ("2. Preprocess", "Skull stripping\n+ CLAHE + Denoise", COLOR_ACCENT),
        ("3. Fusion", "ResNet50 (2048D)\n+ Handcrafted (98D)", COLOR_AMBER),
        ("4. QIEO Solver", "Feature selection\n1,072 optimal dims", COLOR_GREEN),
        ("5. Classifier", "SVM / XGBoost\nFine-tuned CNN", COLOR_RED),
    ]
    for i, (name, desc, color) in enumerate(steps):
        x = Inches(0.6 + i * 1.88)
        add_card(s_arch, x, Inches(2.0), Inches(1.75), Inches(2.2), border_color=color)
        add_textbox(s_arch, x + Inches(0.05), Inches(2.1), Inches(1.65), Inches(0.35), name, font_size=12, color=color, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(s_arch, x + Inches(0.05), Inches(2.5), Inches(1.65), Inches(1.5), desc, font_size=10, color=COLOR_MUTED, alignment=PP_ALIGN.CENTER)
        
    # Detail boxes below
    add_card(s_arch, Inches(0.6), Inches(4.4), Inches(4.5), Inches(2.2))
    add_textbox(s_arch, Inches(0.8), Inches(4.5), Inches(4.1), Inches(0.3), "🔬 Preprocessing details", font_size=14, color=COLOR_ACCENT, bold=True)
    add_bullet_list(s_arch, Inches(0.8), Inches(4.9), Inches(4.1), Inches(1.6), [
        "• Skull Stripping: Otsu thresholding + contour masking.",
        "• Contrast Enhancement: CLAHE (Clip limit = 2.0).",
        "• Denoising: Gaussian filter smoothing.",
        "• Standard scaling: Intensity mapped to [0.0, 1.0]."
    ], font_size=11, color=COLOR_MUTED)

    add_card(s_arch, Inches(5.4), Inches(4.4), Inches(4.5), Inches(2.2))
    add_textbox(s_arch, Inches(5.6), Inches(4.5), Inches(4.1), Inches(0.3), "🧬 Feature Engineering details", font_size=14, color=COLOR_AMBER, bold=True)
    add_bullet_list(s_arch, Inches(5.6), Inches(4.9), Inches(4.1), Inches(1.6), [
        "• Deep CNN features: 2,048 dimensions from ResNet-50.",
        "• Texture GLCM: Contrast, Homogeneity, Energy (20D).",
        "• Contours HOG: Gradient orientation histograms (36D).",
        "• Wavelet DWT: Haar multi-scale frequency bands (32D).",
        "• Intensity: Histograms stats (mean, std, variance) (10D)."
    ], font_size=11, color=COLOR_MUTED)

    print("Created Slide 4: System Architecture")

    # ============================================================
    # SLIDE 5 — QIEO SOLVER
    # ============================================================
    s_qieo = create_cloned_slide()
    add_slide_title(s_qieo, "Core Innovation: The QIEO Solver")
    
    # Left Column Card (Math)
    add_card(s_qieo, Inches(0.6), Inches(2.0), Inches(4.5), Inches(4.6))
    add_textbox(s_qieo, Inches(0.8), Inches(2.1), Inches(4.1), Inches(0.4), "⚛️ Mathematical Formulation", font_size=16, color=COLOR_ACCENT, bold=True)
    add_bullet_list(s_qieo, Inches(0.8), Inches(2.6), Inches(4.1), Inches(3.8), [
        "• Each parameter is represented as a Q-bit state vector:",
        "      q_j = [α_j, β_j]ᵀ   where |α_j|² + |β_j|² = 1",
        "• |α_j|² = probability of collapse to 0 (unselected).",
        "• |β_j|² = probability of collapse to 1 (selected).",
        "• Superposition is maintained during population generation.",
        "• Quantum Rotation Gate applied to update Q-bit amplitudes:",
        "      [α_j', β_j']ᵀ = [cos(θ_j) -sin(θ_j); sin(θ_j) cos(θ_j)] [α_j, β_j]ᵀ",
        "• Rotation angle θ_j is determined dynamically by comparing current solution fitness to the global best."
    ], font_size=11, color=COLOR_MUTED)

    # Right Column Card (Applications)
    add_card(s_qieo, Inches(5.4), Inches(2.0), Inches(4.5), Inches(4.6))
    add_textbox(s_qieo, Inches(5.6), Inches(2.1), Inches(4.1), Inches(0.4), "⚙️ QIEO Applications in QT-2.21", font_size=16, color=COLOR_BLUE, bold=True)
    add_bullet_list(s_qieo, Inches(5.6), Inches(2.6), Inches(4.1), Inches(3.8), [
        "• QIEO Feature Selection:",
        "  - Search space: 2,146 binary features.",
        "  - Fitness function balances training accuracy & mask size.",
        "  - Compresses features by ~50% (selecting 1,072 optimal features).",
        "  - 3-fold CV subset evaluation saves training overhead.",
        "",
        "• QIEO Hyperparameter Optimization:",
        "  - Maps 7 continuous XGBoost parameters to a 56-bit representation.",
        "  - Continuous parameters decoded using binary-to-float scaling.",
        "  - Converges to global minima faster than grid or random search."
    ], font_size=11, color=COLOR_MUTED)

    print("Created Slide 5: QIEO Solver")

    # ============================================================
    # SLIDE 6 — MODEL CONFIGURATIONS
    # ============================================================
    s_conf = create_cloned_slide()
    add_slide_title(s_conf, "Model Configurations: Classical vs. Quantum")
    
    # Table layout
    headers = ["Configuration", "Feature Selector", "Hyperparam Search", "Classifier", "Training Time"]
    row_data = [
        ("Baseline A", "Mutual Info (100)", "GridSearchCV", "SVM (RBF)", "~20 min"),
        ("Baseline B", "SelectKBest (100)", "Optuna TPE (25 trials)", "XGBoost", "~30 min"),
        ("Baseline C", "None (End-to-End)", "Transfer Learning", "ResNet-50", "~3.5 hours"),
        ("Quantum A", "QIEO Mask (1,072)", "GridSearchCV", "SVM (RBF)", "~45 min"),
        ("Quantum B", "QIEO Mask (1,072)", "QIEO 56-bit Binary", "XGBoost", "~35 min"),
        ("Quantum C", "QIEO Informed", "QIEO Informed", "ResNet-50 CNN", "Shared"),
    ]
    
    col_widths = [Inches(1.8), Inches(2.0), Inches(2.1), Inches(1.8), Inches(1.5)]
    left_start = Inches(0.6)
    top_start = Inches(2.0)
    
    # Render table header
    for j, h in enumerate(headers):
        x = left_start + sum(col_widths[:j], Emu(0))
        add_card(s_conf, x, top_start, col_widths[j], Inches(0.45), bg_color=COLOR_ACCENT, border_color=None)
        add_textbox(s_conf, x + Inches(0.05), top_start + Inches(0.08), col_widths[j] - Inches(0.1), Inches(0.35),
                    h, font_size=11, color=RGBColor(255,255,255), bold=True, alignment=PP_ALIGN.CENTER)
                    
    # Render rows
    for i, row in enumerate(row_data):
        y = top_start + Inches(0.55 + i * 0.65)
        is_quantum = "Quantum" in row[0]
        bg = RGBColor(0xFA, 0xF5, 0xFF) if is_quantum else COLOR_CARD_BG
        border = COLOR_ACCENT if is_quantum else COLOR_CARD_BORD
        
        for j, val in enumerate(row):
            x = left_start + sum(col_widths[:j], Emu(0))
            add_card(s_conf, x, y, col_widths[j], Inches(0.6), bg_color=bg, border_color=border)
            color = COLOR_ACCENT if (is_quantum and j == 0) else COLOR_DARK
            add_textbox(s_conf, x + Inches(0.05), y + Inches(0.12), col_widths[j] - Inches(0.1), Inches(0.4),
                        val, font_size=10.5, color=color, bold=(j==0 or is_quantum), alignment=PP_ALIGN.CENTER)
                        
    # Bottom footnote
    add_textbox(s_conf, Inches(0.6), Inches(6.3), Inches(9.2), Inches(0.4),
                "⚛️ QIEO configurations achieve 50% feature compression and optimal hyperparameter alignment in under 45 minutes.",
                font_size=11, color=COLOR_ACCENT, bold=True)

    print("Created Slide 6: Model Configurations")

    # ============================================================
    # SLIDE 7 — RESULTS
    # ============================================================
    s_res = create_cloned_slide()
    add_slide_title(s_res, "Performance Evaluation Results")
    
    # Results Table
    res_headers = ["Config", "Accuracy", "Sensitivity", "Specificity", "Precision", "F1-Macro", "AUC-ROC"]
    res_rows = [
        ("Baseline A", "87.00%", "87.00%", "95.67%", "87.48%", "86.99%", "97.59%"),
        ("Baseline B", "87.00%", "87.00%", "95.67%", "87.17%", "86.96%", "97.46%"),
        ("Baseline C", "95.33%", "95.33%", "98.44%", "95.46%", "95.34%", "99.41%"),
        ("Quantum A", "89.84%", "89.84%", "96.61%", "89.98%", "89.83%", "98.62%"),
        ("Quantum B", "89.29%", "89.29%", "96.43%", "89.50%", "89.28%", "97.99%"),
        ("Quantum C", "95.33%", "95.33%", "98.44%", "95.46%", "95.34%", "99.41%"),
    ]
    
    r_widths = [Inches(1.6), Inches(1.2), Inches(1.3), Inches(1.3), Inches(1.3), Inches(1.3), Inches(1.2)]
    r_left = Inches(0.6)
    r_top = Inches(2.0)
    
    # Header
    for j, h in enumerate(res_headers):
        x = r_left + sum(r_widths[:j], Emu(0))
        add_card(s_res, x, r_top, r_widths[j], Inches(0.45), bg_color=COLOR_ACCENT, border_color=None)
        add_textbox(s_res, x + Inches(0.02), r_top + Inches(0.08), r_widths[j] - Inches(0.04), Inches(0.35),
                    h, font_size=11, color=RGBColor(255,255,255), bold=True, alignment=PP_ALIGN.CENTER)
                    
    # Rows
    for i, row in enumerate(res_rows):
        y = r_top + Inches(0.55 + i * 0.6)
        is_quantum = "Quantum" in row[0]
        bg = RGBColor(0xFA, 0xF5, 0xFF) if is_quantum else COLOR_CARD_BG
        border = COLOR_ACCENT if is_quantum else COLOR_CARD_BORD
        
        for j, val in enumerate(row):
            x = r_left + sum(r_widths[:j], Emu(0))
            add_card(s_res, x, y, r_widths[j], Inches(0.55), bg_color=bg, border_color=border)
            color = COLOR_GREEN if (row[0] in ("Baseline C", "Quantum C") and j > 0) else (COLOR_ACCENT if (is_quantum and j == 0) else COLOR_DARK)
            add_textbox(s_res, x + Inches(0.02), y + Inches(0.1), r_widths[j] - Inches(0.04), Inches(0.35),
                        val, font_size=11, color=color, bold=(j==0 or is_quantum or row[0] in ("Baseline C", "Quantum C")), alignment=PP_ALIGN.CENTER)
                        
    # Bottom analysis boxes
    add_card(s_res, Inches(0.6), Inches(5.8), Inches(4.5), Inches(1.1), border_color=COLOR_GREEN)
    add_textbox(s_res, Inches(0.8), Inches(5.85), Inches(4.1), Inches(0.3), "📈 QIEO Performance Improvements", font_size=13, color=COLOR_GREEN, bold=True)
    add_bullet_list(s_res, Inches(0.8), Inches(6.15), Inches(4.1), Inches(0.7), [
        "• Quantum A vs. Baseline A: +3.26% Accuracy & F1-Score.",
        "• Quantum B vs. Baseline B: +2.63% Accuracy & F1-Score."
    ], font_size=10.5, color=COLOR_MUTED)

    add_card(s_res, Inches(5.3), Inches(5.8), Inches(4.6), Inches(1.1), border_color=COLOR_ACCENT)
    add_textbox(s_res, Inches(5.5), Inches(5.85), Inches(4.2), Inches(0.3), "⚖️ Statistical Significance (McNemar's)", font_size=13, color=COLOR_ACCENT, bold=True)
    add_bullet_list(s_res, Inches(5.5), Inches(6.15), Inches(4.2), Inches(0.7), [
        "• Baseline A vs. Quantum A: p = 0.0037  ✅ Statistically Significant.",
        "• Baseline B vs. Quantum B: p = 0.0110  ✅ Statistically Significant."
    ], font_size=10.5, color=COLOR_MUTED)

    print("Created Slide 7: Performance Evaluation")

    # ============================================================
    # SLIDE 8 — PER-CLASS DIAGNOSTIC PERFORMANCE
    # ============================================================
    s_class = create_cloned_slide()
    add_slide_title(s_class, "Per-Class Diagnostic Performance (Quantum A)")
    
    card_data = [
        ("Glioma", "82.78%", "97.79%", COLOR_RED),
        ("Meningioma", "86.81%", "97.54%", COLOR_AMBER),
        ("Pituitary", "95.24%", "99.52%", COLOR_ACCENT),
        ("No Tumor (Healthy)", "94.51%", "99.62%", COLOR_GREEN)
    ]
    
    # 4 horizontal cards
    for i, (name, sens, auc, color) in enumerate(card_data):
        x = Inches(0.6 + i * 2.33)
        add_card(s_class, x, Inches(2.0), Inches(2.2), Inches(4.6), border_color=color)
        add_textbox(s_class, x + Inches(0.1), Inches(2.2), Inches(2.0), Inches(0.4), name, font_size=14, color=color, bold=True, alignment=PP_ALIGN.CENTER)
        
        # Sensitivity Section
        add_textbox(s_class, x + Inches(0.1), Inches(2.8), Inches(2.0), Inches(0.3), "Sensitivity", font_size=11, color=COLOR_MUTED, alignment=PP_ALIGN.CENTER)
        add_textbox(s_class, x + Inches(0.1), Inches(3.1), Inches(2.0), Inches(0.7), sens, font_size=28, color=COLOR_DARK, bold=True, alignment=PP_ALIGN.CENTER)
        
        # Divider line
        add_card(s_class, x + Inches(0.3), Inches(4.0), Inches(1.6), Inches(0.02), bg_color=COLOR_CARD_BORD, border_color=None)
        
        # AUC-ROC Section
        add_textbox(s_class, x + Inches(0.1), Inches(4.3), Inches(2.0), Inches(0.3), "AUC-ROC", font_size=11, color=COLOR_MUTED, alignment=PP_ALIGN.CENTER)
        add_textbox(s_class, x + Inches(0.1), Inches(4.6), Inches(2.0), Inches(0.7), auc, font_size=28, color=COLOR_DARK, bold=True, alignment=PP_ALIGN.CENTER)
        
        # Analysis Note inside card
        add_textbox(s_class, x + Inches(0.1), Inches(5.5), Inches(2.0), Inches(0.9),
                    "High sensitivity ensures minimal false negatives.", font_size=9.5, color=COLOR_MUTED, alignment=PP_ALIGN.CENTER)

    print("Created Slide 8: Per-class Performance")

    # ============================================================
    # SLIDE 9 — EXPLAINABILITY
    # ============================================================
    s_exp = create_cloned_slide()
    add_slide_title(s_exp, "Explainability & Clinical Trust")
    
    # Left Card (Visual Explainability)
    add_card(s_exp, Inches(0.6), Inches(2.0), Inches(4.5), Inches(4.6))
    add_textbox(s_exp, Inches(0.8), Inches(2.1), Inches(4.1), Inches(0.4), "🔥 Visual Explainability", font_size=16, color=COLOR_RED, bold=True)
    add_bullet_list(s_exp, Inches(0.8), Inches(2.6), Inches(4.1), Inches(3.8), [
        "• Grad-CAM Heatmaps:",
        "  - Maps activations of ResNet-50 layer4 ('base.7').",
        "  - Displays macro-level regions of interest.",
        "  - Visually verifies the CNN targets pathological tissue.",
        "",
        "• Vanilla Saliency Maps:",
        "  - Computes the absolute gradient of predicted class score.",
        "  - Shows edge borders and micro-level tissue alignments."
    ], font_size=12, color=COLOR_MUTED)

    # Right Card (Clinical Reasoning Engine)
    add_card(s_exp, Inches(5.4), Inches(2.0), Inches(4.5), Inches(4.6))
    add_textbox(s_exp, Inches(5.6), Inches(2.1), Inches(4.1), Inches(0.4), "🩺 Clinical Reasoning Engine", font_size=16, color=COLOR_GREEN, bold=True)
    add_bullet_list(s_exp, Inches(5.6), Inches(2.6), Inches(4.1), Inches(3.8), [
        "• Real-time automated diagnosis reasoning including:",
        "  - Texture (GLCM): contrast, homogeneity, energy.",
        "  - Intensity: mean brightness, skewness, variance.",
        "  - Frequencies (Wavelet): DWT subband power levels.",
        "  - Shape (HOG): directional gradient boundaries.",
        "• Grad-CAM Localization descriptor (e.g. 'superior-left').",
        "• Differential Diagnoses reporting runner-up classification.",
        "• Automated radiologist warnings for narrow margin (<15%) decisions."
    ], font_size=12, color=COLOR_MUTED)

    print("Created Slide 9: Explainability")

    # ============================================================
    # SLIDE 10 — LIVE DEMO DASHBOARD
    # ============================================================
    s_demo = create_cloned_slide()
    add_slide_title(s_demo, "Live Interactive Diagnostic Dashboard")
    
    # 4 Grid Features
    features_data = [
        ("📁 MRI Upload & Normalizer", "Drag-and-drop file upload. Previews the skull-stripped and CLAHE contrast-enhanced images in real-time.", COLOR_BLUE),
        ("🧠 Model Selector Panel", "Instantly toggle between 4 classifiers (Quantum A/B vs. Baseline A/B) to see output comparison.", COLOR_ACCENT),
        ("📊 Spatial Attention Maps", "Displays live side-by-side Grad-CAM heatmaps and Saliency maps generated from the PyTorch backend.", COLOR_AMBER),
        ("🩺 Diagnostic Report Sheet", "Outputs the automated clinical reasoning text detailing features, focal location, and differential outcomes.", COLOR_GREEN),
    ]
    for i, (title, desc, color) in enumerate(features_data):
        row = i // 2
        col = i % 2
        x = Inches(0.6 + col * 4.8)
        y = Inches(2.0 + row * 2.4)
        
        add_card(s_demo, x, y, Inches(4.4), Inches(2.1), border_color=color)
        add_textbox(s_demo, x + Inches(0.2), y + Inches(0.15), Inches(4.0), Inches(0.35), title, font_size=14, color=color, bold=True)
        add_textbox(s_demo, x + Inches(0.2), y + Inches(0.6), Inches(4.0), Inches(1.3), desc, font_size=11, color=COLOR_MUTED)

    print("Created Slide 10: Live Demo Dashboard")

    # ============================================================
    # SLIDE 11 — LIMITATIONS & FUTURE SCOPE
    # ============================================================
    s_limit = create_cloned_slide()
    add_slide_title(s_limit, "Clinical Limitations & Future Directions")
    
    # Left Column Card (Limitations)
    add_card(s_limit, Inches(0.6), Inches(2.0), Inches(4.5), Inches(4.6))
    add_textbox(s_limit, Inches(0.8), Inches(2.1), Inches(4.1), Inches(0.4), "⚠️ Clinical Limitations", font_size=16, color=COLOR_AMBER, bold=True)
    add_bullet_list(s_limit, Inches(0.8), Inches(2.6), Inches(4.1), Inches(3.8), [
        "1. Resolution Downsampling:",
        "   Slices resized to 224×224 pixels. Micro-metastases or very tiny lesions could be lost during resizing.",
        "",
        "2. Single Modality (MRI Only):",
        "   Model trained exclusively on MRI. Unusable for CT scan intensity distributions.",
        "",
        "3. Single-Center Bias:",
        "   Trained on a single curated dataset. Robustness to scanner field strengths (1.5T vs. 3T) is unproven."
    ], font_size=11, color=COLOR_MUTED)

    # Right Column Card (Future Directions)
    add_card(s_limit, Inches(5.4), Inches(2.0), Inches(4.5), Inches(4.6))
    add_textbox(s_limit, Inches(5.6), Inches(2.1), Inches(4.1), Inches(0.4), "🚀 Future Scope", font_size=16, color=COLOR_GREEN, bold=True)
    add_bullet_list(s_limit, Inches(5.6), Inches(2.6), Inches(4.1), Inches(3.8), [
        "1. 3D Volumetric Segmentation:",
        "   Extend inputs from 2D slices to 3D volumetric MRI scans using 3D-UNet structures.",
        "",
        "2. Real Quantum Hardware Execution:",
        "   Implement QIEO optimization on physical NISQ computers (IBM Quantum / Rigetti via Qiskit/PennyLane).",
        "",
        "3. Multi-Modal Medical Fusion:",
        "   Fuse scan features with text clinical reports and genomic biomarkers for holistic diagnostics."
    ], font_size=11, color=COLOR_MUTED)

    print("Created Slide 11: Limitations & Future Scope")

    # ============================================================
    # SLIDE 12 — CONCLUSION (Duplicate and customize)
    # ============================================================
    s_conc = create_cloned_slide()
    add_slide_title(s_conc, "Conclusion & Key Takeaways")
    
    takeaways = [
        ("✅", "Quantum A (QIEO + SVM) achieves 89.84% accuracy (+3.26% over classical baseline)."),
        ("✅", "ResNet-50 CNN achieves 95.33% accuracy with 99.41% AUC-ROC."),
        ("✅", "McNemar's test (p = 0.0037) confirms quantum performance gain is statistically significant."),
        ("✅", "Explainability module integrates Grad-CAM heatmaps and Saliency maps for medical auditability."),
        ("✅", "Full end-to-end prototype: Preprocessing → Feature Fusion → QIEO → Classification → Dashboard.")
    ]
    for i, (icon, text) in enumerate(takeaways):
        y = Inches(2.0 + i * 0.8)
        add_card(s_conc, Inches(0.6), y, Inches(9.2), Inches(0.65))
        add_textbox(s_conc, Inches(0.8), y + Inches(0.12), Inches(8.8), Inches(0.45),
                    f"{icon}   {text}", font_size=13, color=COLOR_DARK, bold=True)
                    
    # Bottom Thank You
    add_textbox(s_conc, Inches(0.6), Inches(6.2), Inches(9.2), Inches(0.5), "Thank You! Questions?", font_size=20, color=COLOR_ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

    print("Created Slide 12: Conclusion")

    # ============================================================
    # SAVE PRESENTATION
    # ============================================================
    # Remove the placeholder Slide 2 (the empty content slide template) so it is clean
    id_list = prs.slides._sldIdLst
    del id_list[1]
    
    prs.save(OUTPUT_PATH)
    print(f"\nSUCCESS: Presentation generated and saved to: {OUTPUT_PATH}")
    print(f"Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
